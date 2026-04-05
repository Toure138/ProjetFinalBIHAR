"""
make_pptx.py — Presentation BIHAR 2026 — style moderne teal avec animations.
Usage : python make_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os, copy

# ─── Palette (style screenshot) ───────────────────────────────────────────────
BG       = RGBColor(0x4D, 0x8F, 0x9C)   # teal principal
BG_DARK  = RGBColor(0x2A, 0x5A, 0x68)   # teal foncé (barre bas, header)
BG_BOX   = RGBColor(0x5E, 0x9E, 0xAB)   # teal boîte
GOLD     = RGBColor(0xC9, 0xA8, 0x4C)   # or (numéros, accents)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RED_LINE = RGBColor(0xC0, 0x39, 0x2B)   # soulignement titre
LIGHT    = RGBColor(0xD6, 0xEC, 0xF0)   # texte clair
GREEN_OK = RGBColor(0x27, 0xAE, 0x60)
ORANGE   = RGBColor(0xE6, 0x7E, 0x22)

OUT = "monitoring/output"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

_cid = [10]  # compteur global d'IDs pour animations

def next_id():
    _cid[0] += 1
    return _cid[0]

# ─── Helpers visuels ──────────────────────────────────────────────────────────

def bg(slide, color=None):
    color = color or BG
    r = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    return r

def bar(slide, color=BG_DARK, height=1.1):
    r = slide.shapes.add_shape(1, 0, Inches(7.5-height), prs.slide_width, Inches(height))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    return r

def bar_top(slide, color=BG_DARK, height=0.95):
    r = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(height))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    return r

def rect(slide, l, t, w, h, fill=BG_BOX, border=WHITE, border_px=1.5):
    r = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if border:
        r.line.color.rgb = border
        r.line.width = Pt(border_px)
    else:
        r.line.fill.background()
    return r

def txt(slide, text, l, t, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.CENTER, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.italic = italic
    return tb

def page_num(slide, n):
    r = slide.shapes.add_shape(1, Inches(12.5), Inches(6.9), Inches(0.83), Inches(0.6))
    r.fill.solid(); r.fill.fore_color.rgb = BG_DARK; r.line.fill.background()
    txt(slide, str(n), 12.5, 6.9, 0.83, 0.6, size=14, bold=True, color=WHITE)

def img(slide, path, l, t, w, h=None):
    if not os.path.exists(path): return None
    if h: return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

def title_with_underline(slide, text, l, t, w, h, size=28, color=WHITE):
    """Titre centré avec soulignement rouge (style screenshot)."""
    tb = txt(slide, text, l, t, w, h, size=size, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Soulignement rouge sous le titre
    ul = slide.shapes.add_shape(1, Inches(l + w*0.25), Inches(t + h - 0.08),
                                Inches(w * 0.5), Inches(0.06))
    ul.fill.solid(); ul.fill.fore_color.rgb = RED_LINE; ul.line.fill.background()
    return tb

# ─── Animations XML ───────────────────────────────────────────────────────────

PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'

def _sp_id(shape):
    return str(shape._element.get('id') or shape.shape_id)

def add_fade_animations(slide, shapes, auto=True, base_delay=300):
    """Ajoute des animations fade-in sur les shapes (auto ou sur clic)."""
    if not shapes:
        return

    def make_par(sp_id, click_idx, delay_ms):
        par_xml = f'''<p:par xmlns:p="{PML}">
  <p:cTn id="{next_id()}" fill="hold">
    <p:stCondLst>
      <p:cond delay="{'indefinite' if not auto and click_idx == 0 else str(delay_ms)}"/>
    </p:stCondLst>
    <p:childTnLst>
      <p:par>
        <p:cTn id="{next_id()}" fill="hold">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst>
            <p:par>
              <p:cTn id="{next_id()}" presetID="10" presetClass="entr" presetSubtype="0"
                     fill="hold" grpId="0"
                     nodeType="{'clickEffect' if not auto else 'withEffect'}">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst>
                  <p:set>
                    <p:cBhvr>
                      <p:cTn id="{next_id()}" dur="1" fill="hold"/>
                      <p:tgtEl><p:spTgt spid="{sp_id}"/></p:tgtEl>
                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                    </p:cBhvr>
                    <p:to><p:strVal val="visible"/></p:to>
                  </p:set>
                  <p:animEffect transition="in" filter="fade">
                    <p:cBhvr>
                      <p:cTn id="{next_id()}" dur="600"/>
                      <p:tgtEl><p:spTgt spid="{sp_id}"/></p:tgtEl>
                    </p:cBhvr>
                  </p:animEffect>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:childTnLst>
  </p:cTn>
</p:par>'''
        return etree.fromstring(par_xml)

    child_pars = []
    for i, sh in enumerate(shapes):
        delay = i * base_delay if auto else 0
        child_pars.append(make_par(_sp_id(sh), i, delay))

    bld_items = "".join(
        f'<p:bldP xmlns:p="{PML}" spid="{_sp_id(sh)}" grpId="{i}"/>'
        for i, sh in enumerate(shapes)
    )

    timing_xml = f'''<p:timing xmlns:p="{PML}">
  <p:tnLst>
    <p:par>
      <p:cTn id="{next_id()}" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="{next_id()}" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrevClick" delay="0"><p:tn/></p:cond>
            </p:prevCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst>{bld_items}</p:bldLst>
</p:timing>'''

    timing_el = etree.fromstring(timing_xml)
    main_seq = timing_el.find('.//' + qn('p:cTn') + '[@nodeType="mainSeq"]/' + qn('p:childTnLst'))
    if main_seq is None:
        main_seq = timing_el.find('.//{%s}childTnLst' % PML)
    for par in child_pars:
        main_seq.append(par)

    # Supprimer timing existant si présent
    sp_tree = slide._element
    for old in sp_tree.findall(qn('p:timing')):
        sp_tree.remove(old)
    sp_tree.append(timing_el)

def add_transition(slide, trans="fade"):
    xml = f'<p:transition xmlns:p="{PML}" spd="fast"><p:{trans}/></p:transition>'
    el = etree.fromstring(xml)
    sp_tree = slide._element
    for old in sp_tree.findall(qn('p:transition')):
        sp_tree.remove(old)
    sp_tree.append(el)

# ─── Helpers slides ───────────────────────────────────────────────────────────

def section_slide(num, title, subtitle=""):
    s = prs.slides.add_slide(BLANK)
    bg(s); bar(s)
    # Grand numéro encadré
    b = rect(s, 5.7, 1.5, 1.9, 1.5, fill=BG_DARK, border=WHITE, border_px=2)
    n = txt(s, num, 5.7, 1.5, 1.9, 1.5, size=48, bold=True, color=GOLD)
    title_with_underline(s, title, 1.5, 3.3, 10.3, 1.0, size=34)
    if subtitle:
        t2 = txt(s, subtitle, 1.5, 4.5, 10.3, 0.7, size=18, color=LIGHT, italic=True)
    add_transition(s)
    add_fade_animations(s, [b, n], auto=True, base_delay=200)
    return s

def content_slide(title, items, pnum=0, img_path=None):
    s = prs.slides.add_slide(BLANK)
    bg(s); bar_top(s); bar(s)
    title_with_underline(s, title, 0.3, 0.05, 12.7, 0.85, size=24)
    add_transition(s)

    txt_w = 7.0 if (img_path and os.path.exists(img_path)) else 12.5
    y = 1.1
    animated = []
    for item in items:
        if item.startswith("###"):
            b = rect(s, 0.3, y, txt_w, 0.5, fill=BG_DARK, border=None)
            t = txt(s, item[3:].strip(), 0.3, y, txt_w, 0.5, size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
            animated += [b, t]; y += 0.55
        elif item.startswith("##"):
            t = txt(s, item[2:].strip(), 0.3, y, txt_w, 0.45, size=14, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
            animated.append(t); y += 0.48
        elif item == "":
            y += 0.18
        else:
            b = rect(s, 0.3, y, txt_w, 0.43, fill=BG_BOX, border=WHITE, border_px=0.5)
            t = txt(s, item, 0.45, y+0.02, txt_w-0.3, 0.4, size=13, color=WHITE, align=PP_ALIGN.LEFT)
            animated += [b, t]; y += 0.46

    if img_path and os.path.exists(img_path):
        p = img(s, img_path, 7.5, 1.05, 5.5, 5.5)
        if p: animated.append(p)

    if pnum: page_num(s, pnum)
    add_fade_animations(s, animated, auto=True, base_delay=150)
    return s

def table_slide(title, headers, rows, note="", pnum=0):
    from pptx.util import Inches, Pt
    s = prs.slides.add_slide(BLANK)
    bg(s); bar_top(s); bar(s)
    title_with_underline(s, title, 0.3, 0.05, 12.7, 0.85, size=24)
    add_transition(s)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl_h  = min(0.55 * n_rows, 5.5)
    tbl = s.shapes.add_table(n_rows, n_cols,
                             Inches(0.4), Inches(1.05),
                             Inches(12.5), Inches(tbl_h)).table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = BG_DARK
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.bold  = True
        p.runs[0].font.size  = Pt(13)
        p.runs[0].font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER
        tbl.columns[j].width = Inches(12.5 / n_cols)

    for i, row in enumerate(rows):
        bg_c = BG_BOX if i % 2 == 0 else RGBColor(0x3E, 0x7A, 0x88)
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg_c
            p = cell.text_frame.paragraphs[0]
            c = GOLD if "★" in str(val) or "retenu" in str(val).lower() else WHITE
            p.runs[0].font.size  = Pt(12)
            p.runs[0].font.color.rgb = c
            p.alignment = PP_ALIGN.CENTER

    if note:
        txt(s, note, 0.4, 1.15 + tbl_h, 12.5, 0.5, size=12, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
    if pnum: page_num(s, pnum)
    return s

def image_slide(title, path, caption="", pnum=0):
    s = prs.slides.add_slide(BLANK)
    bg(s); bar_top(s); bar(s)
    title_with_underline(s, title, 0.3, 0.05, 12.7, 0.85, size=24)
    add_transition(s)
    animated = []
    if os.path.exists(path):
        p = img(s, path, 0.5, 1.0, 12.3, 5.4)
        if p: animated.append(p)
    if caption:
        t = txt(s, caption, 0.4, 6.55, 12.5, 0.5, size=11, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)
        animated.append(t)
    if pnum: page_num(s, pnum)
    add_fade_animations(s, animated, auto=True, base_delay=200)
    return s

def two_img_slide(title, p1, p2, c1="", c2="", pnum=0):
    s = prs.slides.add_slide(BLANK)
    bg(s); bar_top(s); bar(s)
    title_with_underline(s, title, 0.3, 0.05, 12.7, 0.85, size=24)
    add_transition(s)
    animated = []
    for path, lx, cap in [(p1, 0.3, c1), (p2, 6.9, c2)]:
        if os.path.exists(path):
            p = img(s, path, lx, 1.0, 6.0, 5.3)
            if p: animated.append(p)
        if cap:
            t = txt(s, cap, lx, 6.4, 6.0, 0.5, size=11, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)
            animated.append(t)
    if pnum: page_num(s, pnum)
    add_fade_animations(s, animated, auto=True, base_delay=300)
    return s

def arch_slide(title, boxes, arrows=None, pnum=0):
    """Slide architecture avec boîtes et flèches."""
    s = prs.slides.add_slide(BLANK)
    bg(s); bar_top(s); bar(s)
    title_with_underline(s, title, 0.3, 0.05, 12.7, 0.85, size=24)
    add_transition(s)
    animated = []

    for box in boxes:
        lx, ty, w, h = box['pos']
        fill = box.get('fill', BG_DARK)
        border = box.get('border', WHITE)
        b = rect(s, lx, ty, w, h, fill=fill, border=border, border_px=1.5)
        animated.append(b)
        if 'title' in box:
            t = txt(s, box['title'], lx, ty+0.05, w, 0.38,
                    size=box.get('title_size', 13), bold=True,
                    color=box.get('title_color', GOLD), align=PP_ALIGN.CENTER)
            animated.append(t)
        if 'body' in box:
            tb = txt(s, box['body'], lx+0.05, ty+0.42, w-0.1,
                     h-0.5, size=box.get('body_size', 11),
                     color=WHITE, align=PP_ALIGN.CENTER)
            animated.append(tb)

    if arrows:
        from pptx.util import Inches
        for arr in arrows:
            x1, y1, x2, y2 = arr
            connector = s.shapes.add_connector(
                2,  # STRAIGHT
                Inches(x1), Inches(y1), Inches(x2), Inches(y2)
            )
            connector.line.color.rgb = WHITE
            connector.line.width = Pt(1.5)

    if pnum: page_num(s, pnum)
    add_fade_animations(s, animated, auto=True, base_delay=200)
    return s

# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

pnum = 1

# ── SLIDE 1 : TITRE ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
bar(s, height=2.2)
bar_top(s, height=0.15, color=GOLD)
animated = []

b1 = rect(s, 0, 0.15, 13.33, 5.15, fill=BG, border=None)
animated.append(b1)

t1 = txt(s, "PROJET FINAL MLOps", 0.5, 0.6, 12.33, 1.1,
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
animated.append(t1)

ul = s.shapes.add_shape(1, Inches(3.5), Inches(1.75), Inches(6.33), Inches(0.07))
ul.fill.solid(); ul.fill.fore_color.rgb = RED_LINE; ul.line.fill.background()
animated.append(ul)

t2 = txt(s, "Prévision de Série Temporelle  •  Classification Multimodale",
         0.5, 1.9, 12.33, 0.7, size=20, color=GOLD, align=PP_ALIGN.CENTER)
animated.append(t2)

for i, (label, val) in enumerate([
    ("Etudiant",   "TOURE Abdoul Aziz"),
    ("Programme",  "MSc BIHAR 2025-2026 — ESTIA"),
    ("Modules",    "M32 Machine Learning II  •  M33 Deep Learning II  •  M27 DevOps"),
]):
    b = rect(s, 1.5, 2.8 + i*0.82, 10.33, 0.65, fill=BG_DARK, border=WHITE, border_px=0.8)
    tl = txt(s, label + " :", 1.7, 2.83 + i*0.82, 2.2, 0.58, size=13, color=GOLD, bold=True, align=PP_ALIGN.LEFT)
    tv = txt(s, val, 3.9, 2.83 + i*0.82, 7.8, 0.58, size=13, color=WHITE, align=PP_ALIGN.LEFT)
    animated += [b, tl, tv]

add_fade_animations(s, animated, auto=True, base_delay=180)
add_transition(s)
pnum += 1

# ── SLIDE 2 : SOMMAIRE ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); bar(s)
add_transition(s)
title_with_underline(s, "Sommaire", 0, 0.1, 13.33, 0.9, size=34)

items = [
    ("01", "Introduction & Contexte"),
    ("02", "Données"),
    ("03", "Classification\nMultimodale"),
    ("04", "Prévision\nSérie Temporelle"),
    ("05", "MLOps / DevOps"),
    ("06", "Résultats &\nConclusion"),
]
animated = []
positions = [(0.4,1.1),(6.9,1.1),(0.4,3.0),(6.9,3.0),(0.4,4.9),(6.9,4.9)]
for (num, title_), (lx, ty) in zip(items, positions):
    b = rect(s, lx, ty, 6.0, 1.6, fill=BG_BOX, border=WHITE, border_px=1.5)
    nb = rect(s, lx, ty, 1.0, 1.6, fill=BG_DARK, border=WHITE, border_px=1.5)
    tn = txt(s, num, lx, ty, 1.0, 1.6, size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tt = txt(s, title_, lx+1.1, ty+0.3, 4.7, 1.0, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    animated += [b, nb, tn, tt]
page_num(s, pnum); pnum += 1
add_fade_animations(s, animated, auto=True, base_delay=120)

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 01 — INTRODUCTION
# ══════════════════════════════════════════════════════════════════════════════
section_slide("01", "Introduction & Contexte")
pnum += 1

content_slide("Contexte du Projet", [
    "###Objectif Global",
    "Construire un pipeline MLOps complet de bout en bout",
    "Deux sous-projets indépendants : Multimodal + Série Temporelle",
    "",
    "###Sous-projet 1 — Classification Multimodale",
    "Prédire plusieurs labels sur des images + descriptions textuelles",
    "Dataset : Kaggle Multi-label Classification Competition 2023",
    "",
    "###Sous-projet 2 — Prévision de Série Temporelle",
    "Prédire la température sur 3 jours (horizon 72h = 24 pas x 3h)",
    "Données : Open-Meteo, Côte d'Ivoire, 2023-2025",
    "",
    "###Partie MLOps",
    "Pipeline entraînement → inférence → API REST → CI/CD → monitoring",
], pnum); pnum += 1

content_slide("Pourquoi ces choix technologiques ?", [
    "###PyTorch — Framework Deep Learning",
    "Flexibilité maximale pour les architectures custom (LSTM, BiLSTM, EfficientNet)",
    "Debuggage facile (eager mode) — meilleur contrôle que TensorFlow sur Windows",
    "",
    "###FastAPI — API REST",
    "Performances élevées (ASGI async) — typage automatique — docs Swagger auto",
    "Natif pour les projets Python ML vs Flask (plus moderne, plus rapide)",
    "",
    "###MLflow — Tracking & Registry",
    "Standard industrie — log params/métriques/artefacts/datasets en 3 lignes",
    "Model Registry intégré — versionning automatique des modèles",
    "",
    "###SQLite — Base de données",
    "Zero-configuration — portable — parfait pour un projet local/batch",
    "Performant pour les volumes de données météo (~10k lignes)",
    "",
    "###GitHub Actions — CI/CD",
    "Natif GitHub — gratuit — GITHUB_TOKEN auto (pas de secrets)",
], pnum); pnum += 1

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 02 — DONNÉES
# ══════════════════════════════════════════════════════════════════════════════
section_slide("02", "Données")
pnum += 1

content_slide("Dataset Multimodal — Kaggle 2023", [
    "###Structure du dataset",
    "~10 000 images (.jpg) + train.csv (ImageID, Labels, Caption)",
    "18 classes de labels (labels 1-19, label 12 absent)",
    "Chaque image peut avoir PLUSIEURS labels simultanément",
    "",
    "###Défi principal : déséquilibre des classes",
    "Classes fréquentes (label 1, 3, 8) : > 2000 occurrences",
    "Classes rares (label 15, 2, 14) : < 500 occurrences — ratio max/min ≈ 10:1",
    "Solution : pos_weight = N_neg/N_pos dans BCEWithLogitsLoss",
    "",
    "###Problème CSV",
    "Captions contiennent des virgules → lignes mal formées",
    "Solution : pd.read_csv(TRAIN_CSV, on_bad_lines='skip')",
    "",
    "###Split : 70% train / 15% val / 15% test",
],
    f"{OUT}/eda_label_distribution.png", pnum); pnum += 1

two_img_slide("EDA Multimodal — Texte & Co-occurrences",
    f"{OUT}/eda_wordcloud.png",
    f"{OUT}/eda_cooccurrence.png",
    "Mots les plus fréquents dans les captions",
    "Labels qui co-apparaissent souvent (heatmap)",
    pnum); pnum += 1

content_slide("Dataset Météo — Série Temporelle", [
    "###Source : API Open-Meteo (gratuite, no auth)",
    "Variable : temperature_2m — Localisation : Côte d'Ivoire (6.82°N, 5.28°W)",
    "Période : 2023-01-01 → 2025-12-31 | ~9 500 observations à 3h",
    "",
    "###Agrégation 3 heures (requis par le sujet)",
    "Valeur à 00h = moyenne(00h, 01h, 02h)",
    "Valeur à 03h = moyenne(03h, 04h, 05h) … etc.",
    "Implémentation : df.resample('3h', label='left').mean()",
    "",
    "###Patterns identifiés par l'EDA",
    "Saisonnalité journalière forte : pic 14h–15h, minimum 06h (amplitude ≈ 8°C)",
    "Saisonnalité annuelle : saison sèche (DJF) vs saison humide (MAM/SON)",
    "Autocorrélation : ρ≈0.97 à lag 1, pic à lag 8 (24h) et lag 56 (7 jours)",
    "Légère tendance haussière : +0.2°C/an sur 2023–2025",
    "",
    "###Split temporel : 70% / 15% / 15% (chronologique)",
    "Fit StandardScaler sur train uniquement — pas de fuite de données",
], pnum); pnum += 1

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 03 — MULTIMODAL
# ══════════════════════════════════════════════════════════════════════════════
section_slide("03", "Classification Multimodale", "Images • Textes • Fusion")
pnum += 1

# Architecture images
arch_slide("Architecture — 3 Modèles Images", [
    {'pos': (0.2,1.1,3.9,1.3), 'title':'Baseline CNN', 'title_size':13,
     'body':'Conv(32)→BN→Pool\nConv(64)→BN→Pool\nConv(128)→BN→Pool\nGAP→Dense(256)→Dense(18)\n622 866 params', 'body_size':10,
     'fill': RGBColor(0x3A,0x70,0x80)},
    {'pos': (4.7,1.1,3.9,1.3), 'title':'EfficientNetB0 FE', 'title_size':13,
     'body':'EfficientNetB0 GELÉ\n(4M params figés)\n+ tête entraînable\nDense(512)→Dense(18)\n23 058 params entr.', 'body_size':10,
     'fill': RGBColor(0x3A,0x70,0x80)},
    {'pos': (9.2,1.1,3.9,1.3), 'title':'EfficientNetB0 FT ★', 'title_size':13, 'title_color': GREEN_OK,
     'body':'3 derniers blocs dégelés\n+ nouvelle tête\n3.18M / 4.03M params\n→ adaptation au domaine', 'body_size':10,
     'fill': RGBColor(0x1A,0x50,0x60)},

    {'pos': (0.2,2.7,12.9,0.5), 'title':'Paramètres communs : Adam lr=1e-3  |  BCEWithLogitsLoss + pos_weight  |  batch=32  |  20 epochs  |  EarlyStopping patience=5', 'title_size':12,
     'fill': BG_DARK},

    {'pos': (0.2,3.4,3.9,2.8), 'title':'Résultats Baseline', 'title_size':12,
     'body':'F1 micro : 0.2873\nF1 macro : 0.2072\nPrécision : 0.1770\nRappel : 0.7621\n\nSous-ajustement\nfeatures CNN limitées', 'body_size':11},
    {'pos': (4.7,3.4,3.9,2.8), 'title':'Résultats FE', 'title_size':12,
     'body':'F1 micro : 0.4544\nF1 macro : 0.3484\nPrécision : 0.3177\nRappel : 0.7974\n\n+18pts vs baseline\nFeatures ImageNet', 'body_size':11},
    {'pos': (9.2,3.4,3.9,2.8), 'title':'Résultats FT ★', 'title_size':12, 'title_color': GREEN_OK,
     'body':'F1 micro : 0.6665\nF1 macro : 0.5553\nPrécision : 0.5748\nRappel : 0.7930\n\n+38pts vs baseline\nMEILLEUR MODELE', 'body_size':11,
     'fill': RGBColor(0x1A,0x50,0x60)},
], pnum=pnum); pnum += 1

two_img_slide("Courbes & Grad-CAM — Classification Images",
    f"{OUT}/img_EfficientNetB0_Fine-tuning_curves.png",
    f"{OUT}/img_gradcam.png",
    "EfficientNetB0 Fine-tuning — courbes d'entraînement",
    "Grad-CAM — zones activées pour chaque prédiction",
    pnum); pnum += 1

# Architecture textes
arch_slide("Architecture — 3 Modèles Textes", [
    {'pos': (0.2,1.1,4.1,1.4), 'title':'TF-IDF + LogReg (baseline)', 'title_size':12,
     'body':'TfidfVectorizer(10 000 features, ngram(1,2))\nOneVsRestClassifier(LogisticRegression)\nRapide, interprétable', 'body_size':10,
     'fill': RGBColor(0x3A,0x70,0x80)},
    {'pos': (4.6,1.1,4.1,1.4), 'title':'TF-IDF + MLP', 'title_size':12,
     'body':'Input(10000)→Dense(512,relu)\n→Dropout→Dense(256)→Dense(18)\nCapture non-linéarités', 'body_size':10,
     'fill': RGBColor(0x3A,0x70,0x80)},
    {'pos': (9.0,1.1,4.1,1.4), 'title':'Bi-LSTM ★', 'title_size':12, 'title_color': GREEN_OK,
     'body':'Embedding(5000, 128)\n→BiLSTM(128)→BiLSTM(64)\n→Dense(128)→Dense(18)\n2.3M params', 'body_size':10,
     'fill': RGBColor(0x1A,0x50,0x60)},

    {'pos': (0.2,2.75,4.1,2.5), 'title':'Résultats Baseline', 'title_size':12,
     'body':'F1 micro : ~0.45\nPas de représentation\nde l\'ordre des mots\nBaseline solide', 'body_size':11},
    {'pos': (4.6,2.75,4.1,2.5), 'title':'Résultats MLP', 'title_size':12,
     'body':'F1 micro : ~0.55\nCapture relations\nnon-linéaires\nentre mots TF-IDF', 'body_size':11},
    {'pos': (9.0,2.75,4.1,2.5), 'title':'Résultats Bi-LSTM ★', 'title_size':12, 'title_color': GREEN_OK,
     'body':'F1 micro : ~0.66\nContexte bidirectionnel\nCapture séquences\nMEILLEUR MODELE', 'body_size':11,
     'fill': RGBColor(0x1A,0x50,0x60)},

    {'pos': (0.2,5.5,12.9,0.6), 'title':'Interprétabilité : LIME — mots les plus influents par prédiction', 'title_size':12,
     'fill': BG_DARK},
], pnum=pnum); pnum += 1

two_img_slide("LIME & Courbes — Classification Textes",
    f"{OUT}/text_bilstm_curves.png",
    f"{OUT}/text_lime_explanation.png",
    "Bi-LSTM — courbes d'entraînement (train / val)",
    "LIME — mots les plus influents pour chaque label prédit",
    pnum); pnum += 1

# Fusion
arch_slide("Architecture — Fusion Multimodale", [
    # Extracteurs
    {'pos': (0.2,1.1,4.0,1.3), 'title':'EfficientNetB0', 'title_size':12,
     'body':'Images 224×224\n→ vecteur 1280-dim', 'body_size':11,
     'fill': RGBColor(0x3A,0x70,0x80)},
    {'pos': (0.2,3.8,4.0,1.3), 'title':'Bi-LSTM', 'title_size':12,
     'body':'Captions tokenisées\n→ vecteur 256-dim', 'body_size':11,
     'fill': RGBColor(0x3A,0x70,0x80)},

    # Early Fusion
    {'pos': (4.6,1.05,4.0,2.1), 'title':'Early Fusion (★)', 'title_size':13, 'title_color': GREEN_OK,
     'body':'Extracteurs GELÉS\nConcat(1536-dim)\n→ Dense(512)\n→ Dense(18)\n5.96M params entr.', 'body_size':10,
     'fill': RGBColor(0x1A,0x50,0x60)},
    # Joint Fusion
    {'pos': (4.6,3.4,4.0,2.1), 'title':'Joint Fusion', 'title_size':13,
     'body':'Extracteurs dégelés\nGate attention:\nw=softmax(Linear(1536,2))\nfused=w0*img+w1*txt\n→ Dense(18)', 'body_size':10},

    # Résultats
    {'pos': (9.0,1.05,4.1,1.0), 'title':'Early Fusion ★', 'title_size':12, 'title_color': GREEN_OK,
     'body':'F1 micro : 0.7383  |  F1 macro : 0.6665\n+13.1% vs image  /  +12.3% vs texte', 'body_size':11,
     'fill': RGBColor(0x1A,0x50,0x60)},
    {'pos': (9.0,2.3,4.1,0.75), 'title':'Joint Fusion', 'title_size':12,
     'body':'F1 micro : 0.6761  |  Attention collapse !\nimage=1.00 / texte=0.00', 'body_size':11},
    {'pos': (9.0,3.4,4.1,0.75), 'title':'Image seule', 'title_size':12,
     'body':'F1 micro : 0.6531', 'body_size':11},
    {'pos': (9.0,4.4,4.1,0.75), 'title':'Texte seul', 'title_size':12,
     'body':'F1 micro : 0.6575', 'body_size':11},

    {'pos': (0.2,5.5,12.9,0.6),
     'title':'Observation : Joint Fusion souffre d\'un attention collapse — image poids=1.0, texte=0.0 → Early Fusion GAGNE',
     'title_size':12, 'title_color': RED_LINE, 'fill': BG_DARK},
], pnum=pnum); pnum += 1

two_img_slide("Résultats Fusion — Comparaison & Attention",
    f"{OUT}/fusion_final_comparison.png",
    f"{OUT}/fusion_attention_weights.png",
    "Comparaison F1 micro des 4 configurations",
    "Poids d'attention Joint Fusion — collapse sur modalité image",
    pnum); pnum += 1

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 04 — SÉRIE TEMPORELLE
# ══════════════════════════════════════════════════════════════════════════════
section_slide("04", "Prévision de Série Temporelle", "LSTM • Prophet • TFT")
pnum += 1

content_slide("Feature Engineering — 49 Variables", [
    "###Pourquoi autant de features pour le LSTM ?",
    "Le LSTM apprend des patterns séquentiels MAIS ne calcule pas de lags explicitement",
    "Fournir les lags et statistiques roulantes donne un avantage décisif (facteur x27 vs TFT)",
    "",
    "###Features cycliques (6) — encodage sin/cos",
    "hour_sin/cos, day_sin/cos, month_sin/cos",
    "Sin/cos évite la discontinuité : heure 23 et 00h sont proches dans l'espace",
    "",
    "###Lag features (26) — autocorrélation",
    "Lags 1→24 (72h), lag 36 (108h), lag 48 (6 jours)",
    "Capturent directement la forte autocorrélation ρ≈0.97",
    "",
    "###Rolling statistics (10) — tendance locale",
    "Moyenne + écart-type sur fenêtres 3h, 6h, 8h (24h), 12h, 24h",
    "",
    "###Différences (3) — taux de variation",
    "diff_1, diff_2, diff_8 — variation relative entre pas",
], pnum); pnum += 1

arch_slide("Architecture LSTM Multi-step", [
    {'pos': (0.3,1.1,2.5,5.5), 'title':'ENTRÉE', 'title_size':12,
     'body':'Séquence\n24 pas × 49 features\n(72h historique)', 'body_size':11,
     'fill': RGBColor(0x2A,0x5A,0x68)},

    {'pos': (3.1,1.1,1.9,1.2), 'title':'LSTM(128)', 'title_size':13, 'title_color': GOLD,
     'body':'Dropout(0.2)', 'body_size':11},
    {'pos': (3.1,2.5,1.9,1.2), 'title':'LSTM(64)', 'title_size':13, 'title_color': GOLD,
     'body':'Dropout(0.2)', 'body_size':11},
    {'pos': (3.1,3.9,1.9,1.2), 'title':'LSTM(32)', 'title_size':13, 'title_color': GOLD,
     'body':'Dropout(0.2)', 'body_size':11},

    {'pos': (5.3,1.1,1.9,1.2), 'title':'Dense(64)', 'title_size':13, 'title_color': LIGHT,
     'body':'ReLU  Dropout(0.2)', 'body_size':11},
    {'pos': (5.3,2.5,1.9,1.2), 'title':'Dense(32)', 'title_size':13, 'title_color': LIGHT,
     'body':'ReLU', 'body_size':11},

    {'pos': (7.5,1.9,1.9,1.5), 'title':'Dense(24)', 'title_size':14, 'title_color': GREEN_OK,
     'body':'Sortie linéaire\n24 pas', 'body_size':12,
     'fill': RGBColor(0x1A,0x50,0x40)},

    {'pos': (9.7,1.1,3.4,5.5), 'title':'PARAMÈTRES', 'title_size':12,
     'body':'Optimizer : Adam\nlr = 1e-3\nLoss : MSELoss\nBatch : 32\nMax epochs : 100\nEarlyStopping : 10\nReduceLROnPlateau :\n  factor=0.5, patience=5\n\nDevice : CPU/CUDA auto', 'body_size':11,
     'fill': RGBColor(0x2A,0x5A,0x68)},

    {'pos': (0.3,6.7,13.0,0.5),
     'title':'Standardscaler fit uniquement sur train — séquences : fenêtre glissante 24 pas → prédit 24 pas suivants',
     'title_size':11, 'fill': BG_DARK},
], pnum=pnum); pnum += 1

table_slide("Résultats — MAE par horizon (jeu de test)",
    ["Horizon", "Prophet\n(baseline stat.)", "LSTM ★\n(modèle prod.)", "TFT\n(exploratoire q50)"],
    [
        ["+3h  (pas 1)", "1.3770", "0.0434", "1.9465"],
        ["+24h (J+1)",   "1.4100", "0.0506", "3.4261"],
        ["+48h (J+2)",   "1.4137", "0.0537", "0.4167"],
        ["+72h (J+3)",   "1.4200", "0.0503", "1.0303"],
        ["Moyenne ★",    "1.4107", "0.0516", "1.5763"],
    ],
    "★ LSTM : MAE ~27× inférieur à Prophet et TFT — Feature engineering = clé du succès",
    pnum); pnum += 1

content_slide("Pourquoi LSTM > TFT ici ?", [
    "###LSTM avec 49 features ingéniérées",
    "Lags explicites + rolling stats → l'information est DIRECTEMENT accessible",
    "Moins de capacité requise pour découvrir les patterns",
    "Entraînement 4 min (CPU) — convergence rapide",
    "",
    "###TFT — Temporal Fusion Transformer",
    "Conçu pour découvrir les patterns AUTOMATIQUEMENT (lags, saisonnalité)",
    "Nécessite plus de données et d'epochs pour rivaliser",
    "Avantages uniques : intervalles de confiance q10-q90 + interprétabilité (VSN)",
    "",
    "###Avantage unique du TFT en production",
    "Variable Selection Networks : quelles features comptent le plus ?",
    "Attention temporelle : quels instants passés influencent la prévision ?",
    "→ Utilisé pour l'analyse et l'interprétabilité, LSTM pour la précision",
    "",
    "###Prophet — Rôle de baseline",
    "Rapide à entraîner (30s), aucun GPU requis",
    "Donne une borne inférieure de performance mesurable",
], pnum); pnum += 1

image_slide("Monitoring — Prédictions LSTM vs Réel",
    f"{OUT}/report_2026-03-25_2026-03-31.png",
    "Comparaison automatique prédictions vs températures observées — MAE calculé par src/monitoring/report.py",
    pnum); pnum += 1

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 05 — MLOps
# ══════════════════════════════════════════════════════════════════════════════
section_slide("05", "MLOps / DevOps", "Architecture • API • CI/CD • Docker • Tests")
pnum += 1

arch_slide("Architecture de l'Application — Flux de Données", [
    # Ligne 1 : collecte
    {'pos': (0.2,1.05,2.8,1.2), 'title':'Open-Meteo API', 'title_size':12,
     'body':'Température 2m\nCôte d\'Ivoire', 'body_size':10, 'fill': RGBColor(0x2A,0x5A,0x68)},
    {'pos': (3.3,1.05,2.8,1.2), 'title':'fetch_data.py', 'title_size':12,
     'body':'Agrégation 3h\nresample().mean()', 'body_size':10},
    {'pos': (6.4,1.05,2.8,1.2), 'title':'weather.db', 'title_size':12,
     'body':'SQLite\nweather_data table', 'body_size':10, 'fill': RGBColor(0x2A,0x5A,0x68)},

    # Ligne 2 : entraînement
    {'pos': (6.4,2.55,2.8,1.2), 'title':'train.py', 'title_size':12,
     'body':'LSTM PyTorch\n49 features', 'body_size':10},
    {'pos': (9.5,2.55,2.8,1.2), 'title':'MLflow', 'title_size':12,
     'body':'params + métriques\n+ artefacts', 'body_size':10, 'fill': RGBColor(0x2A,0x5A,0x68)},
    {'pos': (9.5,1.05,2.8,1.2), 'title':'model/registry/', 'title_size':12,
     'body':'lstm_model.pth\nscaler.pkl', 'body_size':10, 'fill': RGBColor(0x2A,0x5A,0x68)},

    # Ligne 3 : inférence
    {'pos': (3.3,2.55,2.8,1.2), 'title':'predict.py', 'title_size':12,
     'body':'Batch inference\n24 pas (72h)', 'body_size':10},
    {'pos': (0.2,2.55,2.8,1.2), 'title':'predictions table', 'title_size':12,
     'body':'model_id, date\nhorizon, predicted', 'body_size':10, 'fill': RGBColor(0x2A,0x5A,0x68)},

    # Ligne 4 : API
    {'pos': (0.2,4.05,5.8,1.2), 'title':'FastAPI — api/main.py', 'title_size':12,
     'body':'GET /predictions  |  GET /predictions/combined\nGET /version  |  GET /health', 'body_size':11,
     'fill': RGBColor(0x1A,0x50,0x60)},
    {'pos': (6.4,4.05,2.8,1.2), 'title':'report.py', 'title_size':12,
     'body':'MAE prédictions\nvs réel → PNG', 'body_size':10},
    {'pos': (9.5,4.05,2.8,1.2), 'title':'monitoring/output/', 'title_size':12,
     'body':'Graphiques\nde suivi', 'body_size':10, 'fill': RGBColor(0x2A,0x5A,0x68)},
], pnum=pnum); pnum += 1

content_slide("API REST FastAPI — 4 Endpoints", [
    "###GET /predictions?date=YYYY-MM-DD",
    "Retourne les 24 prédictions (72h) générées pour une date donnée",
    'Réponse : {date, model_id, n_steps:24, predictions:[{horizon_step, target_time, predicted}]}',
    "",
    "###GET /predictions/combined?start=...&end=...",
    "Joint prédictions + températures réelles (monitoring post-hoc)",
    'Calcule et retourne : MAE, biais, erreur par point',
    "",
    "###GET /version",
    "software_version : 0.0.0 en local — commit SHA injecté en CI/CD",
    "model_version : MLflow run_id du dernier modèle entraîné",
    "",
    "###GET /health",
    "Retourne {status: ok} — utilisé par Docker et CI/CD",
    "",
    "###Logging automatique sur chaque endpoint",
    "Format : timestamp | level | GET /path → status (latence ms)",
], pnum); pnum += 1

arch_slide("Pipeline CI/CD — GitHub Actions", [
    {'pos': (0.2,1.05,3.9,4.5), 'title':'Job 1 : test', 'title_size':14, 'title_color': GOLD,
     'body':'Déclenché : push + PR\n\n1. Checkout code\n2. Setup Python 3.11\n   + pip cache\n3. pip install -r requirements.txt\n4. mkdir data/ model/\n   monitoring/\n5. pytest tests/ -v\n   → 7 tests\n\nRésultat : 7 passed\n   3.09s', 'body_size':11},
    {'pos': (4.4,1.05,3.9,4.5), 'title':'Job 2 : build', 'title_size':14, 'title_color': GOLD,
     'body':'Déclenché : push seul\nDépend de : test\n\n1. Checkout code\n2. Login ghcr.io\n   GITHUB_TOKEN auto\n   (pas de secret)\n3. docker build\n   --build-arg\n   GIT_COMMIT_ID\n   =${{github.sha}}\n4. Push :latest + :sha\n   → ghcr.io/[owner]/\n   temperature-forecast', 'body_size':11},
    {'pos': (8.6,1.05,4.5,4.5), 'title':'Job 3 : integration-test', 'title_size':14, 'title_color': GOLD,
     'body':'Déclenché : push seul\nDépend de : build\n\n1. Pull image GHCR\n2. docker run -d\n   -p 8000:8000\n3. Wait 30s /health=ok\n4. Test /health → "ok"\n5. Test /version\n   → contient SHA\n6. Test /predictions\n   date=2099 → 404\n7. docker stop (always)', 'body_size':11},

    {'pos': (0.2,5.8,12.9,0.5),
     'title':'Securite : aucun secret en clair — GITHUB_TOKEN genere automatiquement par GitHub Actions',
     'title_size':12, 'fill': BG_DARK},
], pnum=pnum); pnum += 1

content_slide("Docker — Couches optimisées", [
    "###Base : python:3.11-slim",
    "Image légère — pas d'Ubuntu complet — réduit la taille finale",
    "",
    "###Layer 2 — Dépendances (cachée si requirements.txt inchangé)",
    "COPY requirements.txt .",
    "RUN pip install --no-cache-dir -r requirements.txt",
    "→ Rebuild UNIQUEMENT si requirements.txt change (pas à chaque push de code)",
    "",
    "###Layer 3 — Code applicatif",
    "COPY src/ api/ model/ data/ ./",
    "→ Rebuild du code sans réinstaller les dépendances",
    "",
    "###Injection de version (CI/CD)",
    "ARG GIT_COMMIT_ID=0.0.0",
    "ENV GIT_COMMIT_ID=${GIT_COMMIT_ID}",
    "→ GET /version retourne le commit SHA exact du build",
    "",
    "###Démarrage",
    'CMD ["uvicorn", "api.main:app", "--host", "0.0.0.0", "--port", "8000"]',
], pnum); pnum += 1

content_slide("Tests — 7 tests (tous passes)", [
    "###Isolation : DB SQLite temporaire (tempfile.mkdtemp())",
    "Patch cfg.DB_PATH avant import de l'app — aucun impact sur la prod",
    "",
    "###Tests unitaires (3)",
    "test_health — GET /health → 200, status=ok",
    "test_version_structure — software_version + model_version présents et de type string",
    "test_predictions_date_format — date invalide → 404 ou 422 (pas de crash)",
    "",
    "###Tests d'intégration (4)",
    "test_predictions_not_found — date inconnue → 404 avec détail",
    "test_predictions_ok — insertion DB → GET /predictions → 200, 24 steps, structure correcte",
    "test_combined_not_found — période vide → 404",
    "test_combined_ok — prédictions + réel insérés → summary.mae calculé",
    "",
    "###Résultat",
    "7 passed in 3.09s — Lancé en CI/CD avant chaque build Docker",
    "Commande : pytest tests/ -v --tb=short",
], pnum); pnum += 1

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 06 — RÉSULTATS
# ══════════════════════════════════════════════════════════════════════════════
section_slide("06", "Résultats & Conclusion")
pnum += 1

table_slide("Récapitulatif — Classification Multimodale",
    ["Sous-tâche", "Modèle retenu", "F1 micro", "F1 macro", "Gain vs baseline"],
    [
        ["Images",  "EfficientNetB0 FT",  "0.6665", "0.5553", "+37.9 pts"],
        ["Textes",  "Bi-LSTM",            "~0.66",  "~0.59",  "+21 pts"],
        ["Fusion ★","Early Fusion",       "0.7383", "0.6665", "+46 pts images"],
    ],
    "★ Early Fusion = meilleure architecture — +13.1% vs image seule — Joint Fusion souffre d'attention collapse",
    pnum); pnum += 1

table_slide("Récapitulatif — Série Temporelle",
    ["Modèle", "MAE global", "RMSE global", "Type", "Avantage"],
    [
        ["Prophet",    "1.4107", "1.8695", "Statistique", "Rapide, baseline"],
        ["LSTM ★",     "0.0516", "0.0652", "Deep Learning", "Précision x27"],
        ["TFT",        "1.5763", "1.8644", "Transformer",   "Intervalles conf."],
    ],
    "★ LSTM déployé en production — Feature engineering = clé (49 variables vs features brutes TFT)",
    pnum); pnum += 1

content_slide("Points Durs & Solutions", [
    "###DLL PyTorch/TensorFlow sur Windows",
    "torch 2.11.0 nightly + TF nécessitent CUDA → OSError c10.dll / pywrap_tensorflow",
    "Solution : pip install torch==2.3.0+cpu --index-url pytorch.org/whl/cpu",
    "",
    "###MLflow URI invalide sur Windows",
    "C:\\... non reconnu comme scheme → UnsupportedModelRegistryStoreURIException",
    "Solution : MLFLOW_DIR.as_uri() génère file:///C:/... reconnu par MLflow",
    "",
    "###pytorch-forecasting API changée v1.0",
    "predict() retourne un namedtuple 5 champs (pas un dict) → ValueError unpack",
    "Solution : raw_preds.output['prediction'] au lieu de raw_preds['prediction']",
    "",
    "###CSV mal formé (captions avec virgules)",
    "ParserError: Expected 3 fields, saw 4 → toutes les lignes avec virgule échouent",
    "Solution : pd.read_csv(csv, on_bad_lines='skip')",
    "",
    "###numpy 2.4.3 incompatible torch 2.3.0",
    "RuntimeError: Numpy is not available → .detach().cpu().numpy() échoue",
    "Solution : pip install numpy==1.26.4",
], pnum); pnum += 1

# ── SLIDE FINAL ───────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
bar(s, height=2.5)
bar_top(s, height=0.15, color=GOLD)
add_transition(s)
animated = []

t1 = txt(s, "Merci de votre attention", 0, 1.2, 13.33, 1.2,
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
ul = s.shapes.add_shape(1, Inches(3.5), Inches(2.45), Inches(6.33), Inches(0.07))
ul.fill.solid(); ul.fill.fore_color.rgb = RED_LINE; ul.line.fill.background()
animated += [t1, ul]

t2 = txt(s, "Demonstration API", 0, 2.8, 13.33, 0.8,
         size=26, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
animated.append(t2)

for i, (label, cmd) in enumerate([
    ("Démarrer l'API :", "uvicorn api.main:app --host 0.0.0.0 --port 8000"),
    ("Documentation :", "http://localhost:8000/docs"),
    ("Générer prédictions :", "python -m src.inference.predict --date 2026-04-01"),
]):
    b = rect(s, 1.5, 3.8 + i*0.78, 10.33, 0.65, fill=BG_DARK, border=WHITE, border_px=0.8)
    tl = txt(s, label, 1.7, 3.83 + i*0.78, 2.5, 0.58, size=12, color=GOLD, bold=True, align=PP_ALIGN.LEFT)
    tv = txt(s, cmd, 4.2, 3.83 + i*0.78, 7.4, 0.58, size=12, color=LIGHT, align=PP_ALIGN.LEFT, italic=True)
    animated += [b, tl, tv]

tf = txt(s, "TOURE Abdoul Aziz  —  MSc BIHAR 2025-2026  —  ESTIA",
         0, 6.9, 13.33, 0.5, size=13, color=LIGHT, align=PP_ALIGN.CENTER)
animated.append(tf)
add_fade_animations(s, animated, auto=True, base_delay=200)

# ══════════════════════════════════════════════════════════════════════════════
prs.save("PRESENTATION_BIHAR2026_TOURE_v2.pptx")
print("OK - Presentation saved: PRESENTATION_BIHAR2026_TOURE_v2.pptx")
print(f"   {len(prs.slides)} slides")
